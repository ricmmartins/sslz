from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Color Palette (Microsoft-inspired dark theme) ---
BG_DARK = RGBColor(0x1B, 0x1B, 0x1B)
BG_CARD = RGBColor(0x2D, 0x2D, 0x2D)
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)
ACCENT_TEAL = RGBColor(0x00, 0xB2, 0x94)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xD1, 0x34, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=18, default_color=WHITE, font_name="Segoe UI"):
    """lines is a list of tuples: (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text, size, color, bold, align = line_data
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size if size else default_size)
        p.font.color.rgb = color if color else default_color
        p.font.bold = bold if bold else False
        p.font.name = font_name
        p.alignment = align if align else PP_ALIGN.LEFT
        p.space_after = Pt(4)
    return txBox

def add_table(slide, left, top, width, height, rows, col_widths=None):
    """rows is list of lists. First row is header."""
    num_rows = len(rows)
    num_cols = len(rows[0])
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_text
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.name = "Segoe UI"
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY

            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ACCENT_BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD if r_idx % 2 == 1 else RGBColor(0x38, 0x38, 0x38)

    return table_shape


# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)

# Accent bar at top
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
            "Startup-Scale Landing Zone", 44, WHITE, True, PP_ALIGN.LEFT)

add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
            "Get Azure right from day one — without the enterprise overhead", 24, LIGHT_GRAY, False, PP_ALIGN.LEFT)

# Divider line
add_shape_bg(slide, Inches(1), Inches(4.2), Inches(2), Inches(0.04), ACCENT_BLUE)

add_textbox(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
            "Ricardo Martins  |  April 2, 2026", 18, MID_GRAY, False, PP_ALIGN.LEFT)

add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
            "startupscalelanding.zone", 18, ACCENT_BLUE, False, PP_ALIGN.LEFT)

# Add notes
slide.notes_slide.notes_text_frame.text = (
    "Open with energy. This is YOUR initiative. You own this.\n\n"
    "\"I want to walk you through something I've been building that I think has real potential "
    "for how we onboard startups onto Azure.\""
)


# ============================================================
# SLIDE 2 — The Gap I Identified
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "The Gap I Identified", 36, WHITE, True)

add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.8),
    [
        ["What Exists Today", "The Problem"],
        ["ALZ (Enterprise-Scale)", "100+ modules, months to understand, built for 10K-seat enterprises"],
        ["ALZ-Bicep", "Still enterprise-scoped, overwhelming for a 10-person startup"],
        ["CAF Terraform Module", "Enterprise-scoped, entering extended support (archived Aug 2026)"],
        ["Nothing", "Most startups just skip governance entirely"],
    ],
    col_widths=[Inches(3.5), Inches(8)]
)

# Highlight box
card = add_shape_bg(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.2), RGBColor(0x00, 0x2B, 0x4E), 0.02)
add_textbox(slide, Inches(1.2), Inches(5.15), Inches(10.5), Inches(1.0),
            "💡 The real competitor isn't ALZ — it's zero governance.", 26, ACCENT_TEAL, True, PP_ALIGN.LEFT)

slide.notes_slide.notes_text_frame.text = (
    "Don't bash existing ALZ — it's great for its audience. Emphasize the GAP, not a flaw.\n\n"
    "Key insight: \"The competition for startups isn't ALZ vs SSLZ. It's governance vs NO governance. "
    "Most startups choose nothing.\"\n\n"
    "Mention the CAF Terraform module entering extended support — this is timely.\n"
    "Pause and let them absorb the table. The visual contrast is powerful."
)


# ============================================================
# SLIDE 3 — What Startups Actually Do
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "What Startups Actually Do Today", 36, WHITE, True)

add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
            "Two paths — both wrong:", 22, LIGHT_GRAY, False)

# Path A card
add_shape_bg(slide, Inches(0.8), Inches(2.1), Inches(5.4), Inches(2.5), BG_CARD, 0.02)
add_multiline_textbox(slide, Inches(1.1), Inches(2.2), Inches(4.9), Inches(2.3), [
    ("Option A: Follow full ALZ", 20, ACCENT_BLUE, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("2 months of \"cloud foundations\" work", 16, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("Team of 5 engineers paralyzed", 16, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("Like buying a commercial kitchen to make breakfast", 14, MID_GRAY, False, PP_ALIGN.LEFT),
])

# Path B card
add_shape_bg(slide, Inches(6.6), Inches(2.1), Inches(5.8), Inches(2.5), BG_CARD, 0.02)
add_multiline_textbox(slide, Inches(6.9), Inches(2.2), Inches(5.3), Inches(2.3), [
    ("Option B: Skip governance", 20, ACCENT_ORANGE, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("One sub, no policies, no budgets, no RBAC", 16, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("Works until:", 16, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("• First enterprise security questionnaire", 14, MID_GRAY, False, PP_ALIGN.LEFT),
    ("• First runaway cost incident", 14, MID_GRAY, False, PP_ALIGN.LEFT),
    ("• First az group delete that hits production", 14, MID_GRAY, False, PP_ALIGN.LEFT),
])

# Bottom message
card = add_shape_bg(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.2), RGBColor(0x00, 0x2B, 0x4E), 0.02)
add_textbox(slide, Inches(1.2), Inches(5.35), Inches(10.5), Inches(1.0),
            "🚀 Startups need a third path.", 28, WHITE, True, PP_ALIGN.LEFT)

slide.notes_slide.notes_text_frame.text = (
    "This is the emotional slide. Make it real with customer anecdotes.\n\n"
    "\"I've seen this pattern dozens of times working with startups — they skip governance, "
    "then six months later they're scrambling when an enterprise customer sends a security questionnaire.\""
)


# ============================================================
# SLIDE 4 — The Solution: SSLZ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "The Solution I Built: SSLZ", 36, WHITE, True)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.6),
            "A deployable, opinionated Azure Landing Zone that ships in < 1 hour", 22, ACCENT_TEAL, False)

features = [
    ("1 Management Group, 2 Subscriptions", "Prod + Non-Prod — that's it. No six-layer hierarchy."),
    ("Security built-in", "Defender, RBAC groups, NSG deny-all defaults, policy enforcement — all automated."),
    ("Cost controls from day one", "Budget alerts at 50/80/100%. Mandatory tagging. No exceptions."),
    ("Explicit graduation path", "Step-by-step guide to migrate to full ALZ when you're ready."),
]

colors = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_GREEN, ACCENT_ORANGE]

for i, (title, desc) in enumerate(features):
    y = 2.1 + i * 1.15
    # Color bar
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(0.08), Inches(0.85), colors[i])
    add_shape_bg(slide, Inches(0.95), Inches(y), Inches(11.2), Inches(0.85), BG_CARD, 0.01)
    add_textbox(slide, Inches(1.2), Inches(y + 0.05), Inches(10.5), Inches(0.4),
                title, 18, WHITE, True)
    add_textbox(slide, Inches(1.2), Inches(y + 0.42), Inches(10.5), Inches(0.4),
                desc, 15, LIGHT_GRAY, False)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
            "Target: Teams of 5–50 engineers  |  Pre-seed to Series A  |  No platform team", 16, MID_GRAY, False)

slide.notes_slide.notes_text_frame.text = (
    "Key emphasis: UNDER 1 HOUR to deploy. Say it clearly.\n\n"
    "\"This isn't a whitepaper or a reference architecture. It's a git clone + az deployment sub create and you're done.\"\n\n"
    "If they ask about scale: \"The graduation guide is explicit — here's when you outgrow it, "
    "and here's exactly how to move to full ALZ.\""
)


# ============================================================
# SLIDE 5 — Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Architecture", 36, WHITE, True)

# MG hierarchy card
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.0), BG_CARD, 0.02)
add_textbox(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.4),
            "Management Group Hierarchy", 18, ACCENT_BLUE, True)
add_multiline_textbox(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(3.0), [
    ("Tenant Root Group", 16, MID_GRAY, False, PP_ALIGN.LEFT),
    ("  └── mg-<yourcompany>", 16, WHITE, True, PP_ALIGN.LEFT),
    ("        ← Policies applied here", 13, ACCENT_TEAL, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("      ├── sub-prod", 16, ACCENT_GREEN, False, PP_ALIGN.LEFT),
    ("           Production workloads", 13, MID_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("      └── sub-nonprod", 16, ACCENT_ORANGE, False, PP_ALIGN.LEFT),
    ("           Dev, staging, QA", 13, MID_GRAY, False, PP_ALIGN.LEFT),
], font_name="Cascadia Code")

# VNet layout card
add_shape_bg(slide, Inches(6.7), Inches(1.5), Inches(5.8), Inches(4.0), BG_CARD, 0.02)
add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.3), Inches(0.4),
            "VNet Layout (per subscription)", 18, ACCENT_BLUE, True)
add_multiline_textbox(slide, Inches(7.0), Inches(2.2), Inches(5.3), Inches(3.0), [
    ("vnet-prod (10.0.0.0/16)", 16, WHITE, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("  snet-aks      10.0.0.0/20    4,091 IPs", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("  snet-app      10.0.16.0/22   1,019 IPs", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("  snet-data     10.0.20.0/22   1,019 IPs", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("  snet-shared   10.0.24.0/24   251 IPs", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
], font_name="Cascadia Code")

# Bottom callout — what's NOT here
add_shape_bg(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0), RGBColor(0x3A, 0x1A, 0x00), 0.02)
add_textbox(slide, Inches(1.2), Inches(5.9), Inches(10.8), Inches(0.8),
            "No hub VNet  •  No Azure Firewall ($900+/mo)  •  No VNet peering  •  Each subscription is self-contained", 18, ACCENT_ORANGE, True)

slide.notes_slide.notes_text_frame.text = (
    "Walk through the diagram. Emphasize what's NOT there:\n"
    "\"No hub VNet — that's $1,500/mo minimum. Startups don't need it until hybrid connectivity.\"\n"
    "\"No Azure Firewall — $900+/mo. NSGs handle 95% of startup networking use cases for free.\"\n\n"
    "If asked about security without a firewall: \"NSGs provide L3/L4 filtering. For most startups with "
    "a single workload in a single region, that's sufficient.\""
)


# ============================================================
# SLIDE 6 — What Ships Out of the Box
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "What Ships Out of the Box", 36, WHITE, True)

add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5),
    [
        ["Component", "What's Deployed"],
        ["Management Groups", "Single MG with 2 subscriptions"],
        ["Azure Policy", "MCSB (audit), required tags, allowed locations, diagnostic settings"],
        ["Networking", "VNet + 4 subnets per subscription, NSGs with deny-all-inbound"],
        ["Monitoring", "Log Analytics workspace, Activity Log forwarding, 90-day retention"],
        ["Security", "Defender CSPM (free) + Servers P2 (prod), security contact alerts"],
        ["Cost Management", "Budget alerts at 50/80/100%, tag enforcement via policy"],
        ["CI/CD", "GitHub Actions for Bicep & Terraform, Workload Identity Federation"],
    ],
    col_widths=[Inches(3), Inches(8.5)]
)

slide.notes_slide.notes_text_frame.text = (
    "Move through this quickly — it's a reference slide.\n\n"
    "Highlight: Defender CSPM is free. Servers P2 on prod only.\n"
    "Workload Identity Federation: No secrets to store, rotate, or commit.\n"
    "Policy in Audit mode, not Deny: \"Don't block legitimate deployments on day one.\""
)


# ============================================================
# SLIDE 7 — Startup Archetypes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Built for Real Startup Archetypes", 36, WHITE, True)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
            "Three production-grade example architectures — full Bicep + Terraform implementations", 18, LIGHT_GRAY, False)

archetypes = [
    ("SaaS Startup", "Container Apps + Azure SQL\nElastic Pool + Redis + Key Vault",
     "Multi-tenant, shared schema,\nscale-to-zero in non-prod", ACCENT_BLUE),
    ("AI Startup", "AKS + GPU Spot Pools +\nAzure OpenAI + Blob + Redis",
     "60–90% GPU savings with Spot,\nKEDA autoscaling", ACCENT_TEAL),
    ("API-First Startup", "App Service + API Management\n+ Cosmos DB + App Insights",
     "Pay-per-call APIM, deployment\nslots, zero-downtime swaps", ACCENT_GREEN),
]

for i, (title, stack, detail, color) in enumerate(archetypes):
    x = 0.8 + i * 4.0
    # Card
    add_shape_bg(slide, Inches(x), Inches(2.0), Inches(3.7), Inches(4.5), BG_CARD, 0.02)
    # Color top bar
    add_shape_bg(slide, Inches(x), Inches(2.0), Inches(3.7), Inches(0.08), color)
    # Title
    add_textbox(slide, Inches(x + 0.3), Inches(2.3), Inches(3.2), Inches(0.5),
                title, 22, color, True)
    # Stack
    add_textbox(slide, Inches(x + 0.3), Inches(3.0), Inches(3.2), Inches(1.2),
                stack, 16, WHITE, False)
    # Detail
    add_textbox(slide, Inches(x + 0.3), Inches(4.5), Inches(3.2), Inches(1.2),
                detail, 14, MID_GRAY, False)

slide.notes_slide.notes_text_frame.text = (
    "\"These aren't toy examples. Each has full Bicep + Terraform, deployment instructions, "
    "and realistic cost estimates.\"\n\n"
    "Pick the most relevant to your audience and go slightly deeper.\n"
    "If time is short, just acknowledge they exist and move on."
)


# ============================================================
# SLIDE 8 — Traction & Blog
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Traction & Thought Leadership", 36, WHITE, True)

# Repo stats card
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.8), BG_CARD, 0.02)
add_textbox(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
            "📦 Repository", 22, ACCENT_BLUE, True)

repo_stats = [
    ("211", "commits in ~6 weeks"),
    ("6+5", "Bicep + Terraform modules"),
    ("5", "GitHub Actions workflows"),
    ("8", "documentation pages"),
    ("3", "startup archetype examples"),
]

for i, (num, desc) in enumerate(repo_stats):
    y = 2.3 + i * 0.7
    add_textbox(slide, Inches(1.3), Inches(y), Inches(1.2), Inches(0.5),
                num, 28, ACCENT_TEAL, True)
    add_textbox(slide, Inches(2.6), Inches(y + 0.05), Inches(3.5), Inches(0.5),
                desc, 16, LIGHT_GRAY, False)

# Blog stats card
add_shape_bg(slide, Inches(6.7), Inches(1.5), Inches(5.8), Inches(4.8), BG_CARD, 0.02)
add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.3), Inches(0.5),
            "📝 Startups at Microsoft Blog", 22, ACCENT_BLUE, True)

blog_stats = [
    ("38", "articles on TechCommunity"),
    ("189K+", "total views across all posts"),
    ("42.5K", "views on top post (From Zero to Hero)"),
]

for i, (num, desc) in enumerate(blog_stats):
    y = 2.3 + i * 0.8
    add_textbox(slide, Inches(7.2), Inches(y), Inches(1.8), Inches(0.5),
                num, 28, ACCENT_TEAL, True)
    add_textbox(slide, Inches(9.1), Inches(y + 0.05), Inches(3.2), Inches(0.5),
                desc, 16, LIGHT_GRAY, False)

# Topics list
add_textbox(slide, Inches(7.2), Inches(4.6), Inches(5), Inches(0.4),
            "Topics covered:", 14, MID_GRAY, True)
add_multiline_textbox(slide, Inches(7.2), Inches(5.0), Inches(5), Inches(1.5), [
    ("Landing Zones • Identity • AKS Networking", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("Monitoring • Cost Optimization • AI Gateway", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("3 posts with 25K+ views each", 13, ACCENT_TEAL, True, PP_ALIGN.LEFT),
])

# Bottom bar
add_shape_bg(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5), RGBColor(0x00, 0x2B, 0x4E), 0.01)
add_textbox(slide, Inches(1.0), Inches(6.63), Inches(11), Inches(0.4),
            "🌐 startupscalelanding.zone  —  Custom domain, Jekyll site, polished & production-ready", 15, ACCENT_TEAL, False)

slide.notes_slide.notes_text_frame.text = (
    "This is your credibility slide. Deliver the numbers with confidence.\n\n"
    "\"211 commits in six weeks. This isn't a side project.\"\n\n"
    "Blog: \"38 articles on TechCommunity. Nearly 190,000 total views.\"\n"
    "Pause after 190K. Let that number register.\n"
    "\"Top 3 posts alone: 42K, 35K, 25K views. We're reaching the right audience.\"\n\n"
    "\"38 articles built the audience. SSLZ gives them something actionable. Founders Hub is the distribution channel.\""
)


# ============================================================
# SLIDE 9 — Design Principles
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Design Principles", 36, WHITE, True)

add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
            "Three rules behind every SSLZ decision:", 20, LIGHT_GRAY, False)

principles = [
    ("1", "Opinionated over flexible",
     "\"It depends\" isn't helpful with 5 engineers and no platform team.\nSSLZ makes the call: two subscriptions, no hub, deny-all NSGs, MCSB in audit mode.",
     ACCENT_BLUE),
    ("2", "Reversible over perfect",
     "Every decision can be changed later. Moving subscriptions between\nmanagement groups is a 10-second operation. Adding a hub = new deployment, not rebuild.",
     ACCENT_TEAL),
    ("3", "Honest about trade-offs",
     "We say \"you'll outgrow this when...\" and \"here's exactly what\nthe next layer costs.\" No false claims of enterprise-grade.",
     ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(principles):
    y = 2.1 + i * 1.7
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(1.45), BG_CARD, 0.02)
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(0.08), Inches(1.45), color)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(y + 0.25), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(2.2), Inches(y + 0.1), Inches(9.5), Inches(0.4),
                title, 22, WHITE, True)
    add_textbox(slide, Inches(2.2), Inches(y + 0.6), Inches(9.5), Inches(0.8),
                desc, 15, LIGHT_GRAY, False)

slide.notes_slide.notes_text_frame.text = (
    "Spend 30 seconds on each principle.\n\n"
    "\"Opinionated\" — Startups don't have time for 'it depends.'\n"
    "\"Reversible\" — Nothing paints you into a corner.\n"
    "\"Honest\" — We don't claim enterprise-grade."
)


# ============================================================
# SLIDE 10 — Founders Hub Integration Vision
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "The Founders Hub Integration Vision", 36, WHITE, True)

# Flow: TODAY
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.8), BG_CARD, 0.02)
add_textbox(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.4),
            "📍 TODAY", 22, ACCENT_RED, True)
add_multiline_textbox(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(3.5), [
    ("Startup joins Founders Hub", 17, WHITE, False, PP_ALIGN.LEFT),
    ("        ↓", 17, MID_GRAY, False, PP_ALIGN.LEFT),
    ("Gets Azure credits + benefits", 17, WHITE, False, PP_ALIGN.LEFT),
    ("        ↓", 17, MID_GRAY, False, PP_ALIGN.LEFT),
    ("No guidance on HOW to set up Azure", 17, ACCENT_RED, True, PP_ALIGN.LEFT),
    ("        ↓", 17, MID_GRAY, False, PP_ALIGN.LEFT),
    ("Burns credits on ungoverned single-sub", 17, ACCENT_ORANGE, False, PP_ALIGN.LEFT),
    ("        ↓", 17, MID_GRAY, False, PP_ALIGN.LEFT),
    ("Security incidents, cost surprises, churn", 17, ACCENT_RED, False, PP_ALIGN.LEFT),
])

# Flow: WITH SSLZ
add_shape_bg(slide, Inches(6.7), Inches(1.6), Inches(5.8), Inches(4.8), BG_CARD, 0.02)
add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.3), Inches(0.4),
            "🚀 WITH SSLZ", 22, ACCENT_GREEN, True)
add_multiline_textbox(slide, Inches(7.0), Inches(2.3), Inches(5.3), Inches(3.5), [
    ("Startup joins Founders Hub", 17, WHITE, False, PP_ALIGN.LEFT),
    ("        ↓", 17, MID_GRAY, False, PP_ALIGN.LEFT),
    ("Gets Azure credits + SSLZ guidance", 17, WHITE, False, PP_ALIGN.LEFT),
    ("        ↓", 17, MID_GRAY, False, PP_ALIGN.LEFT),
    ("Deploys secure foundation in 1 hour", 17, ACCENT_GREEN, True, PP_ALIGN.LEFT),
    ("        ↓", 17, MID_GRAY, False, PP_ALIGN.LEFT),
    ("Security + cost controls from day one", 17, ACCENT_TEAL, False, PP_ALIGN.LEFT),
    ("        ↓", 17, MID_GRAY, False, PP_ALIGN.LEFT),
    ("Graduates to full ALZ when ready", 17, ACCENT_BLUE, False, PP_ALIGN.LEFT),
])

# Arrow between
add_textbox(slide, Inches(5.8), Inches(3.5), Inches(1.5), Inches(0.6),
            "→", 48, ACCENT_TEAL, True, PP_ALIGN.CENTER)

slide.notes_slide.notes_text_frame.text = (
    "THIS IS THE ASK SLIDE. Slow down.\n\n"
    "\"When a startup joins Founders Hub today, they get credits and benefits. "
    "But there's no guidance on HOW to set up Azure properly.\"\n\n"
    "\"SSLZ fills that gap. It's the Day 0 playbook.\"\n\n"
    "Be specific: \"I'm proposing we link SSLZ from Founders Hub onboarding as a recommended resource. "
    "Not mandatory — recommended.\"\n\n"
    "\"This is complementary to ALZ, not competitive. It's the on-ramp.\""
)


# ============================================================
# SLIDE 11 — Why This Matters for the Org
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Why This Matters for the Org", 36, WHITE, True)

impacts = [
    ("Startup retention on Azure", "Better first experience → longer engagement on the platform",
     "📈", ACCENT_BLUE),
    ("Support ticket reduction", "Governance from day 1 → fewer \"my bill exploded\" escalations",
     "🎫", ACCENT_TEAL),
    ("Enterprise readiness", "Startups graduate to ALZ → become enterprise Azure customers",
     "🏢", ACCENT_GREEN),
    ("Content flywheel", "Blog series drives awareness → SSLZ drives adoption → rinse & repeat",
     "🔄", ACCENT_ORANGE),
]

for i, (title, desc, emoji, color) in enumerate(impacts):
    y = 1.6 + i * 1.3
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(1.05), BG_CARD, 0.02)
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(0.08), Inches(1.05), color)
    add_textbox(slide, Inches(1.2), Inches(y + 0.02), Inches(0.6), Inches(0.5),
                emoji, 28, WHITE, False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.9), Inches(y + 0.05), Inches(4), Inches(0.45),
                title, 20, WHITE, True)
    add_textbox(slide, Inches(1.9), Inches(y + 0.5), Inches(9.5), Inches(0.5),
                desc, 16, LIGHT_GRAY, False)

slide.notes_slide.notes_text_frame.text = (
    "Connect SSLZ to business outcomes THEY care about.\n\n"
    "\"Better first experience means startups stay on Azure longer.\"\n"
    "\"Governance from day one means fewer support tickets.\"\n"
    "\"Startups that start with SSLZ become full ALZ enterprise customers.\"\n"
    "\"Blog + SSLZ = flywheel.\""
)


# ============================================================
# SLIDE 12 — The Ask
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "What I'm Asking For", 36, WHITE, True)

asks = [
    ("Endorsement", "Position SSLZ as a Founders Hub recommended resource"),
    ("Collaboration", "Intro to Founders Hub team for onboarding integration"),
    ("Feedback", "Graduation path alignment with ALZ team's vision"),
    ("Visibility", "Share with startup-facing field teams as a starting point"),
]

icons = ["✅", "🤝", "💬", "📢"]

for i, (title, desc) in enumerate(asks):
    y = 1.8 + i * 1.3
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(1.05), BG_CARD, 0.02)
    add_textbox(slide, Inches(1.2), Inches(y + 0.05), Inches(0.6), Inches(0.5),
                icons[i], 28, WHITE, False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.0), Inches(y + 0.05), Inches(3.5), Inches(0.45),
                title, 22, ACCENT_BLUE, True)
    add_textbox(slide, Inches(2.0), Inches(y + 0.52), Inches(9.5), Inches(0.5),
                desc, 16, LIGHT_GRAY, False)

slide.notes_slide.notes_text_frame.text = (
    "Be direct. You're presenting a plan and requesting support.\n\n"
    "1. Endorsement: \"I'd like your backing.\"\n"
    "2. Collaboration: \"I need an intro to the Founders Hub team.\"\n"
    "3. Feedback: \"I want alignment with the ALZ team.\"\n"
    "4. Visibility: \"If you think this is valuable, share it.\""
)


# ============================================================
# SLIDE 13 — Timeline & Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "Timeline & Next Steps", 36, WHITE, True)

steps = [
    ("Identify gap in startup Azure onboarding", "✅ Done", ACCENT_GREEN),
    ("Build SSLZ (Bicep + Terraform + docs)", "✅ Done", ACCENT_GREEN),
    ("Publish on TechCommunity blog", "✅ Done (Mar 12)", ACCENT_GREEN),
    ("Launch startupscalelanding.zone", "✅ Done", ACCENT_GREEN),
    ("Present to team (Ravi, Bhaskar, Amit)", "📍 Today", ACCENT_BLUE),
    ("Founders Hub onboarding integration", "🔜 Next", ACCENT_ORANGE),
    ("Field team enablement & distribution", "🔜 Planned", ACCENT_ORANGE),
    ("Community feedback loop & iteration", "🔄 Ongoing", ACCENT_TEAL),
]

for i, (action, status, color) in enumerate(steps):
    y = 1.5 + i * 0.7
    bg_color = BG_CARD if i % 2 == 0 else RGBColor(0x38, 0x38, 0x38)
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(0.6), bg_color, 0.01)
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(0.06), Inches(0.6), color)
    add_textbox(slide, Inches(1.2), Inches(y + 0.08), Inches(7.5), Inches(0.4),
                action, 16, WHITE if i < 5 else LIGHT_GRAY, False)
    add_textbox(slide, Inches(9.0), Inches(y + 0.08), Inches(3), Inches(0.4),
                status, 16, color, True)

slide.notes_slide.notes_text_frame.text = (
    "Walk through quickly. The checkmarks build credibility — you've already done the hard work.\n\n"
    "\"The solution is built and published. Now I need org support to drive integration.\""
)


# ============================================================
# SLIDE 14 — Close
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_BLUE)

# Divider
add_shape_bg(slide, Inches(4.5), Inches(1.5), Inches(4), Inches(0.04), ACCENT_BLUE)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.0),
            "I identified the gap.\nI built the solution.\nNow let's drive it into the org.", 36, WHITE, True, PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(4.5), Inches(3.8), Inches(4), Inches(0.04), ACCENT_BLUE)

add_multiline_textbox(slide, Inches(1), Inches(4.3), Inches(11.333), Inches(2.0), [
    ("🔗  startupscalelanding.zone", 22, ACCENT_BLUE, False, PP_ALIGN.CENTER),
    ("📦  github.com/ricmmartins/sslz", 22, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("📝  Startups at Microsoft Blog — 38 articles, ~190K views", 20, MID_GRAY, False, PP_ALIGN.CENTER),
])

# Quote
add_shape_bg(slide, Inches(2), Inches(5.8), Inches(9), Inches(0.9), BG_CARD, 0.02)
add_textbox(slide, Inches(2.3), Inches(5.9), Inches(8.5), Inches(0.7),
            "\"For startups, the alternative isn't ALZ — it's usually no governance at all.\"", 20, ACCENT_TEAL, True, PP_ALIGN.CENTER)

slide.notes_slide.notes_text_frame.text = (
    "End strong and circle back to ownership.\n\n"
    "\"I identified this gap through my work with startups. I built the solution. "
    "I published the blog series. Now I'm asking for your help to drive it into the org.\"\n\n"
    "Close with the quote. Then: \"What questions do you have?\""
)


# ============================================================
# SAVE
# ============================================================
output_dir = r"C:\Users\ricardomac\.copilot\session-state\7e089d29-6c99-4eb4-89c5-0d207c905d45\files"
output_path = os.path.join(output_dir, "SSLZ-Presentation.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
